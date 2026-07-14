import fitz
import zipfile
import os
from io import BytesIO
from PIL import Image
from django.http import JsonResponse
from django.views.decorators.http import require_POST
from django.views.decorators.csrf import csrf_exempt
from apps.pdf_tools.utils import save_uploaded_file, get_output_path, media_url, cleanup_file, validate_pdf, safe_int
from apps.pdf_tools.mongo_db import save_job


@csrf_exempt
@require_POST
def pdf_to_jpg(request):
    f = request.FILES.get('file')
    dpi = safe_int(request.POST.get('dpi', 150), default=150, min_val=72, max_val=600)
    fmt = request.POST.get('format', 'jpg').lower()

    if not f:
        return JsonResponse({'error': 'No file uploaded.'}, status=400)

    saved_path, _ = save_uploaded_file(f)
    img_paths = []

    try:
        doc = fitz.open(saved_path)
        matrix = fitz.Matrix(dpi/72, dpi/72)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix)
            img_path, img_name = get_output_path(f'.{fmt}', f'page_{i+1}')
            if fmt == 'jpg':
                img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
                img.save(img_path, 'JPEG', quality=95)
            else:
                pix.save(img_path)
            img_paths.append((img_path, img_name))

        doc.close()

        if len(img_paths) == 1:
            save_job('pdf_to_jpg', [f.name], [img_paths[0][1]])
            return JsonResponse({'download_url': media_url(img_paths[0][1]), 'filename': img_paths[0][1]})

        zip_path, zip_name = get_output_path('.zip', 'pdf_images')
        with zipfile.ZipFile(zip_path, 'w') as zf:
            for ip, iname in img_paths:
                zf.write(ip, iname)

        save_job('pdf_to_jpg', [f.name], [zip_name])
        return JsonResponse({'download_url': media_url(zip_name), 'filename': zip_name, 'pages': len(img_paths)})
    except ValueError as e:
        return JsonResponse({'error': str(e)}, status=400)
    except Exception:
        return JsonResponse({'error': 'Conversion failed. Ensure the file is a valid PDF.'}, status=500)
    finally:
        cleanup_file(saved_path)


@csrf_exempt
@require_POST
def pdf_to_word(request):
    f = request.FILES.get('file')
    if not f:
        return JsonResponse({'error': 'No file uploaded.'}, status=400)

    saved_path, _ = save_uploaded_file(f)
    out_path, out_name = get_output_path('.docx', 'pdf_to_word')

    try:
        validate_pdf(saved_path, f.name)
        from pdf2docx import Converter
        cv = Converter(saved_path)
        cv.convert(out_path)
        cv.close()
        save_job('pdf_to_word', [f.name], [out_name])
        return JsonResponse({'download_url': media_url(out_name), 'filename': out_name})
    except ValueError as e:
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'error': f'Conversion failed: {e}'}, status=500)
    finally:
        cleanup_file(saved_path)


@csrf_exempt
@require_POST
def pdf_to_pptx(request):
    """Convert PDF pages to PowerPoint (each page = 1 slide as image)."""
    f = request.FILES.get('file')
    if not f:
        return JsonResponse({'error': 'No file uploaded.'}, status=400)

    saved_path, _ = save_uploaded_file(f)
    out_path, out_name = get_output_path('.pptx', 'pdf_to_pptx')
    img_paths = []

    try:
        validate_pdf(saved_path, f.name)
        from pptx import Presentation
        from pptx.util import Inches

        doc = fitz.open(saved_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_path, img_name = get_output_path('.png', f'slide_{i+1}')
            pix.save(img_path)
            img_paths.append(img_path)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

        doc.close()
        prs.save(out_path)

        save_job('pdf_to_pptx', [f.name], [out_name])
        return JsonResponse({'download_url': media_url(out_name), 'filename': out_name})
    except ValueError as e:
        return JsonResponse({'error': str(e)}, status=400)
    except Exception:
        return JsonResponse({'error': 'Conversion failed. Ensure the file is a valid PDF.'}, status=500)
    finally:
        cleanup_file(saved_path)
        for ip in img_paths:
            cleanup_file(ip)


@csrf_exempt
@require_POST
def pdf_to_excel(request):
    """Extract tables from PDF and save to Excel."""
    f = request.FILES.get('file')
    if not f:
        return JsonResponse({'error': 'No file uploaded.'}, status=400)

    saved_path, _ = save_uploaded_file(f)
    out_path, out_name = get_output_path('.xlsx', 'pdf_to_excel')

    try:
        validate_pdf(saved_path, f.name)
        import pdfplumber
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        sheet_idx = 0

        with pdfplumber.open(saved_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    for tbl_idx, table in enumerate(tables):
                        ws = wb.create_sheet(title=f'P{page_num+1}_T{tbl_idx+1}')
                        header_fill = PatternFill('solid', fgColor='4F81BD')
                        for row_idx, row in enumerate(table):
                            for col_idx, cell in enumerate(row):
                                c = ws.cell(row=row_idx+1, column=col_idx+1, value=cell or '')
                                if row_idx == 0:
                                    c.font = Font(bold=True, color='FFFFFF')
                                    c.fill = header_fill
                                c.alignment = Alignment(wrap_text=True)
                        sheet_idx += 1
                else:
                    # No table — put raw text
                    ws = wb.create_sheet(title=f'Page_{page_num+1}')
                    text = page.extract_text() or ''
                    for i, line in enumerate(text.split('\n')):
                        ws.cell(row=i+1, column=1, value=line)
                    sheet_idx += 1

        if not wb.sheetnames:
            wb.create_sheet('Sheet1')

        wb.save(out_path)
        save_job('pdf_to_excel', [f.name], [out_name])
        return JsonResponse({'download_url': media_url(out_name), 'filename': out_name})
    except ValueError as e:
        return JsonResponse({'error': str(e)}, status=400)
    except Exception:
        return JsonResponse({'error': 'Conversion failed. Ensure the file is a valid PDF.'}, status=500)
    finally:
        cleanup_file(saved_path)

